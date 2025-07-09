from typing import Dict, List, Optional
import pandas as pd
import numpy as np
from datetime import datetime
import logging
from pathlib import Path
import yaml
import json
import matplotlib.pyplot as plt
import seaborn as sns
from jinja2 import Template
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import nbformat as nbf
from nbformat.v4 import new_notebook, new_markdown_cell, new_code_cell

class ReportAgent:
    def __init__(self, config_path: Optional[str] = None, llm_config: Optional[Dict] = None):
        """
        Initialize the Report Generation Agent.
        
        Args:
            config_path (str, optional): Path to configuration file
            llm_config (dict, optional): Configuration for LLM
        """
        self.logger = logging.getLogger(__name__)
        self.setup_logging()
        self.config = self.load_config(config_path)
        self.llm_config = llm_config or {}
        
    def setup_logging(self):
        """Configure logging."""
        log_dir = Path('logs')
        log_dir.mkdir(exist_ok=True)
        
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_dir / f'report_{datetime.now().strftime("%Y%m%d_%H%M%S")}.log'),
                logging.StreamHandler()
            ]
        )
    
    def load_config(self, config_path: Optional[str]) -> Dict:
        """Load configuration from YAML file."""
        default_config = {
            'report': {
                'template_dir': 'templates',
                'output_dir': 'reports',
                'sections': [
                    'executive_summary',
                    'data_overview',
                    'forecast_analysis',
                    'market_insights',
                    'simulation_results',
                    'recommendations'
                ]
            },
            'visualization': {
                'style': 'seaborn',
                'context': 'notebook',
                'palette': 'deep',
                'figsize': (12, 6)
            }
        }
        
        if config_path and Path(config_path).exists():
            with open(config_path, 'r') as f:
                return {**default_config, **yaml.safe_load(f)}
        return default_config
    
    def generate_executive_summary(self, forecast_results: Dict, 
                                 market_insights: Dict,
                                 simulation_results: Dict) -> str:
        """
        Generate executive summary focusing on unusual events and news analysis.
        
        Args:
            forecast_results (Dict): Forecasting results
            market_insights (Dict): Market analysis results
            simulation_results (Dict): Simulation results
            
        Returns:
            str: Executive summary
        """
        summary = []
        
        # Title
        summary.append("# Unusual Events and Market Analysis Report")
        summary.append(f"\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        
        # Unusual Events Analysis
        summary.append("## Unusual Events Analysis")
        if 'anomalies' in forecast_results:
            summary.append("\n### Detected Anomalies")
            for anomaly in forecast_results['anomalies']:
                summary.append(f"- Date: {anomaly['date']}")
                summary.append(f"  - Value: {anomaly['value']}")
                summary.append(f"  - Deviation: {anomaly['deviation']:.2f}%")
                summary.append(f"  - Impact: {anomaly['impact']}\n")
        
        # Market News Analysis
        summary.append("## Market News Analysis")
        if 'articles' in market_insights:
            summary.append("\n### Relevant News Articles")
            for article in market_insights['articles']:
                summary.append(f"- Title: {article['title']}")
                summary.append(f"  - Date: {article['published_time']}")
                summary.append(f"  - Sentiment: {article['sentiment']['sentiment']}")
                summary.append(f"  - Key Points: {article['content'][:200]}...\n")
        
        # Recommendations
        summary.append("## Recommendations")
        if 'recommendations' in market_insights:
            for rec in market_insights['recommendations']:
                summary.append(f"- {rec}")
        
        return "\n".join(summary)
    
    def create_interactive_dashboard(self, forecast_results: Dict,
                                   market_insights: Dict,
                                   simulation_results: Dict) -> go.Figure:
        """
        Create interactive dashboard using Plotly.
        
        Args:
            forecast_results (Dict): Forecasting results
            market_insights (Dict): Market analysis results
            simulation_results (Dict): Simulation results
            
        Returns:
            go.Figure: Interactive dashboard
        """
        # Create subplot figure
        fig = make_subplots(
            rows=2, cols=2,
            subplot_titles=(
                "Forecast Results",
                "Market Trends",
                "Scenario Analysis",
                "Feature Importance"
            )
        )
        
        # Add forecast plot
        fig.add_trace(
            go.Scatter(
                x=forecast_results['dates'],
                y=forecast_results['predictions'],
                name='Forecast',
                line=dict(color='blue')
            ),
            row=1, col=1
        )
        
        # Add market trends
        for trend in market_insights['trends']:
            fig.add_trace(
                go.Scatter(
                    x=trend['dates'],
                    y=trend['values'],
                    name=trend['name'],
                    line=dict(color='green')
                ),
                row=1, col=2
            )
        
        # Add scenario analysis
        for scenario in simulation_results['scenarios']:
            fig.add_trace(
                go.Scatter(
                    x=[scenario['feature_change']],
                    y=[scenario['predictions']['mean']],
                    name=f"Scenario {scenario['scenario_id']}",
                    error_y=dict(
                        type='data',
                        array=[scenario['predictions']['std']],
                        visible=True
                    )
                ),
                row=2, col=1
            )
        
        # Add feature importance
        importance_df = pd.DataFrame(simulation_results['feature_importance'])
        fig.add_trace(
            go.Bar(
                x=importance_df['importance'],
                y=importance_df['feature'],
                orientation='h',
                name='Feature Importance'
            ),
            row=2, col=2
        )
        
        # Update layout
        fig.update_layout(
            height=800,
            showlegend=True,
            title_text="Interactive Analysis Dashboard"
        )
        
        return fig
    
    def generate_notebook(self, forecast_results: Dict,
                         market_insights: Dict,
                         simulation_results: Dict) -> str:
        """
        Generate Jupyter notebook with analysis.
        
        Args:
            forecast_results (Dict): Forecasting results
            market_insights (Dict): Market analysis results
            simulation_results (Dict): Simulation results
            
        Returns:
            str: Path to generated notebook
        """
        nb = new_notebook()
        
        # Add title
        nb.cells.append(new_markdown_cell("# Analysis Report"))
        
        # Add executive summary
        nb.cells.append(new_markdown_cell("## Executive Summary"))
        nb.cells.append(new_markdown_cell(
            self.generate_executive_summary(forecast_results, market_insights, simulation_results)
        ))
        
        # Add forecast analysis
        nb.cells.append(new_markdown_cell("## Forecast Analysis"))
        nb.cells.append(new_code_cell(
            f"""
            import plotly.graph_objects as go
            
            fig = go.Figure()
            fig.add_trace(go.Scatter(
                x={forecast_results['dates']},
                y={forecast_results['predictions']},
                name='Forecast'
            ))
            fig.show()
            """
        ))
        
        # Add market insights
        nb.cells.append(new_markdown_cell("## Market Insights"))
        for insight in market_insights['key_insights']:
            nb.cells.append(new_markdown_cell(f"- {insight}"))
        
        # Add simulation results
        nb.cells.append(new_markdown_cell("## Simulation Results"))
        nb.cells.append(new_code_cell(
            f"""
            import pandas as pd
            
            scenarios_df = pd.DataFrame({simulation_results['scenarios']})
            scenarios_df
            """
        ))
        
        # Save notebook
        output_dir = Path(self.config['report']['output_dir'])
        output_dir.mkdir(parents=True, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        notebook_path = output_dir / f'analysis_report_{timestamp}.ipynb'
        
        with open(notebook_path, 'w') as f:
            nbf.write(nb, f)
        
        return str(notebook_path)
    
    def generate_presentation(self, forecast_results: Dict,
                            market_insights: Dict,
                            simulation_results: Dict) -> str:
        """
        Generate PowerPoint presentation focusing on unusual events and news analysis.
        
        Args:
            forecast_results (Dict): Forecasting results
            market_insights (Dict): Market analysis results
            simulation_results (Dict): Simulation results
            
        Returns:
            str: Path to generated presentation
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            title.text = "Unusual Events and Market Analysis"
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
            
            # Unusual Events slide
            events_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = events_slide.shapes.title
            content = events_slide.placeholders[1]
            title.text = "Unusual Events Analysis"
            
            events_text = []
            if 'anomalies' in forecast_results:
                for anomaly in forecast_results['anomalies']:
                    events_text.append(f"• Date: {anomaly['date']}")
                    events_text.append(f"  Value: {anomaly['value']}")
                    events_text.append(f"  Deviation: {anomaly['deviation']:.2f}%")
                    events_text.append(f"  Impact: {anomaly['impact']}\n")
            content.text = "\n".join(events_text)
            
            # Market News slide
            news_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = news_slide.shapes.title
            content = news_slide.placeholders[1]
            title.text = "Market News Analysis"
            
            news_text = []
            if 'articles' in market_insights:
                for article in market_insights['articles']:
                    news_text.append(f"• {article['title']}")
                    news_text.append(f"  Date: {article['published_time']}")
                    news_text.append(f"  Sentiment: {article['sentiment']['sentiment']}")
                    news_text.append(f"  Key Points: {article['content'][:100]}...\n")
            content.text = "\n".join(news_text)
            
            # Recommendations slide
            rec_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = rec_slide.shapes.title
            content = rec_slide.placeholders[1]
            title.text = "Recommendations"
            
            if 'recommendations' in market_insights:
                content.text = "\n".join([f"• {rec}" for rec in market_insights['recommendations']])
            
            # Save presentation
            output_dir = Path(self.config['report']['output_dir'])
            output_dir.mkdir(parents=True, exist_ok=True)
            presentation_path = output_dir / f'analysis_presentation_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
            prs.save(str(presentation_path))
            
            return str(presentation_path)
            
        except Exception as e:
            self.logger.error(f"Error generating presentation: {str(e)}")
            raise
    
    def generate_report(self, forecast_results: Dict,
                       market_insights: Dict,
                       simulation_results: Dict) -> Dict:
        """
        Generate comprehensive report focusing on unusual events and news analysis.
        
        Args:
            forecast_results (Dict): Forecasting results
            market_insights (Dict): Market analysis results
            simulation_results (Dict): Simulation results
            
        Returns:
            Dict: Generated reports
        """
        try:
            # Create output directory
            output_dir = Path(self.config['report']['output_dir'])
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Generate executive summary focusing on unusual events
            summary = self.generate_executive_summary(forecast_results, market_insights, simulation_results)
            
            # Create presentation
            presentation_path = self.generate_presentation(forecast_results, market_insights, simulation_results)
            
            # Save summary
            summary_path = output_dir / 'executive_summary.md'
            with open(summary_path, 'w') as f:
                f.write(summary)
            
            return {
                'generated_reports': {
                    'executive_summary': str(summary_path),
                    'presentation': presentation_path
                }
            }
            
        except Exception as e:
            self.logger.error(f"Error generating report: {str(e)}")
            raise 