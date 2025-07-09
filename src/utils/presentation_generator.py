import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import json
from datetime import datetime

class PresentationGenerator:
    def __init__(self, output_dir: str = 'output'):
        """Initialize the presentation generator."""
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
    def generate_presentation(self, data: pd.DataFrame, forecast_results: dict, analysis_results: dict):
        """Generate a PowerPoint presentation with market data summary."""
        prs = Presentation()
        
        # Title Slide
        self._add_title_slide(prs, "NIQ Hackfest - Market Analysis Report")
        
        # Market Overview
        self._add_market_overview(prs, data)
        
        # Key Findings
        self._add_key_findings(prs, analysis_results)
        
        # Forecast Analysis
        self._add_forecast_analysis(prs, forecast_results)
        
        # Trends and Patterns
        self._add_trends_patterns(prs, data)
        
        # Recommendations
        self._add_recommendations(prs, analysis_results)
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(self.output_dir, f'market_analysis_{timestamp}.pptx')
        prs.save(output_path)
        return output_path
        
    def _add_title_slide(self, prs, title):
        """Add title slide."""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
    def _add_market_overview(self, prs, data):
        """Add market overview slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = "Market Overview"
        
        # Add market statistics
        stats = {
            "Total Sales": f"${data['sales_value'].sum():,.2f}",
            "Average Sales": f"${data['sales_value'].mean():,.2f}",
            "Total Volume": f"{data['sales_volume'].sum():,.0f} units",
            "Number of Products": f"{data['SingleBuyProductItemId'].nunique():,}",
            "Number of Clusters": f"{data['ClusterId'].nunique():,}",
            "Average Outlets": f"{data['num_outlets'].mean():,.1f}"
        }
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for key, value in stats.items():
            p = tf.add_paragraph()
            p.text = f"{key}: {value}"
            p.font.size = Pt(18)
            p.font.bold = True
            
    def _add_key_findings(self, prs, analysis_results):
        """Add key findings slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = "Key Findings"
        
        # Add key findings from analysis
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        findings = analysis_results.get('findings', [])
        for finding in findings:
            p = tf.add_paragraph()
            p.text = f"• {finding}"
            p.font.size = Pt(16)
            
    def _add_forecast_analysis(self, prs, forecast_results):
        """Add forecast analysis slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = "Forecast Analysis"
        
        # Add forecast metrics
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        metrics = forecast_results.get('metrics', {})
        for metric, value in metrics.items():
            p = tf.add_paragraph()
            p.text = f"{metric.upper()}: {value:.2f}"
            p.font.size = Pt(18)
            p.font.bold = True
            
    def _add_trends_patterns(self, prs, data):
        """Add trends and patterns slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = "Trends and Patterns"
        
        # Add trend analysis
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Calculate trends
        monthly_sales = data.groupby(data.index.month)['sales_value'].mean()
        monthly_volume = data.groupby(data.index.month)['sales_volume'].mean()
        
        sales_trend = "Increasing" if monthly_sales.iloc[-1] > monthly_sales.iloc[0] else "Decreasing"
        volume_trend = "Increasing" if monthly_volume.iloc[-1] > monthly_volume.iloc[0] else "Decreasing"
        
        # Add trends
        p = tf.add_paragraph()
        p.text = f"Sales Trend: {sales_trend}"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"Volume Trend: {volume_trend}"
        p.font.size = Pt(18)
        p.font.bold = True
        
        # Add cluster analysis
        cluster_stats = data.groupby('ClusterId').agg({
            'sales_value': 'sum',
            'sales_volume': 'sum',
            'num_outlets': 'mean'
        }).round(2)
        
        p = tf.add_paragraph()
        p.text = "\nCluster Analysis:"
        p.font.size = Pt(16)
        p.font.bold = True
        
        for cluster_id, stats in cluster_stats.iterrows():
            p = tf.add_paragraph()
            p.text = f"Cluster {cluster_id}: Sales ${stats['sales_value']:,.2f}, Volume {stats['sales_volume']:,.0f} units, Avg Outlets {stats['num_outlets']:,.1f}"
            p.font.size = Pt(14)
        
    def _add_recommendations(self, prs, analysis_results):
        """Add recommendations slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = "Recommendations"
        
        # Add recommendations
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        recommendations = analysis_results.get('recommendations', [])
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = f"• {rec}"
            p.font.size = Pt(16) 