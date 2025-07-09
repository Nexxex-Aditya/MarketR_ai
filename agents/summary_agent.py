# This is the summary agent which will combines outputs from all agents and generate a natural language report + PowerPoint Slides.

# Create ppt_generator.py utility for slide generation.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any, Optional
import json
from datetime import datetime
import pandas as pd
import numpy as np
from utils.llm_utils import call_llm
import logging
from utils.ppt_utils import create_ppt_report # Import the utility function
import os
from pathlib import Path

logger = logging.getLogger(__name__)

def create_price_trend_chart(prs: Presentation, data: List[Dict[str, Any]]) -> None:
    """Create a slide with price trend analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Price Trend Analysis"
    
    # Convert to DataFrame for easier analysis
    df = pd.DataFrame(data)
    df['Date'] = pd.to_datetime(df['Date'])
    df = df.sort_values('Date')
    
    # Calculate moving averages
    df['MA7'] = df['Price'].rolling(window=7).mean()
    df['MA30'] = df['Price'].rolling(window=30).mean()
    
    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = [d.strftime('%Y-%m-%d') for d in df['Date']]
    chart_data.add_series('Price', df['Price'].tolist())
    chart_data.add_series('7-day MA', df['MA7'].tolist())
    chart_data.add_series('30-day MA', df['MA30'].tolist())
    
    # Create chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    # Format chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.plots[0].has_data_labels = False
    
    # Add trend analysis text
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(1)
    )
    text_frame = textbox.text_frame
    
    # Calculate trend
    price_change = ((df['Price'].iloc[-1] - df['Price'].iloc[0]) / df['Price'].iloc[0]) * 100
    trend_text = f"Price Change: {price_change:.1f}% over the period"
    
    p = text_frame.add_paragraph()
    p.text = trend_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)

def create_impact_distribution_chart(prs: Presentation, simulations: List[Dict[str, Any]]) -> None:
    """Create a slide with impact distribution analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Impact Distribution Analysis"
    
    # Count impact levels
    impact_counts = {}
    for sim in simulations:
        impact = sim['impact_assessment']
        impact_counts[impact] = impact_counts.get(impact, 0) + 1
    
    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = list(impact_counts.keys())
    chart_data.add_series('Count', list(impact_counts.values()))
    
    # Create chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # Format chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = True

def create_news_sentiment_chart(prs: Presentation, news_items: List[Dict[str, Any]]) -> None:
    """Create a slide with news sentiment analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "News Sentiment Analysis"
    
    # Simple sentiment analysis based on keywords
    sentiments = {
        'positive': ['increase', 'growth', 'opportunity', 'positive'],
        'negative': ['decrease', 'decline', 'risk', 'negative'],
        'neutral': ['stable', 'maintain', 'unchanged']
    }
    
    sentiment_counts = {'positive': 0, 'negative': 0, 'neutral': 0}
    
    for article in news_items:
        text = article['title'] + ' ' + article['summary']
        text = text.lower()
        
        if any(word in text for word in sentiments['positive']):
            sentiment_counts['positive'] += 1
        elif any(word in text for word in sentiments['negative']):
            sentiment_counts['negative'] += 1
        else:
            sentiment_counts['neutral'] += 1
    
    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = list(sentiment_counts.keys())
    chart_data.add_series('Count', list(sentiment_counts.values()))
    
    # Create chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    # Format chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.plots[0].has_data_labels = True

def create_news_slide(prs: Presentation, news_results: Dict[str, Any]) -> None:
    """Create a slide with news highlights based on news_results from NewsAgent."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market News Highlights"
    
    textbox = slide.shapes.placeholders[1].text_frame
    textbox.clear()
    
    articles = news_results.get('articles', [])
    if not articles:
        p = textbox.add_paragraph()
        p.text = "No relevant news articles found."
        p.font.size = Pt(14)
        return

    for article in articles[:3]:  # Show top 3 articles
        p = textbox.add_paragraph()
        p.text = f"📰 {article.get('title', 'N/A')}"
        p.font.bold = True
        p.font.size = Pt(14)
        
        p = textbox.add_paragraph()
        p.text = f"Source: {article.get('source', 'N/A')} | Published: {article.get('date', 'N/A')}"
        p.font.italic = True
        p.font.size = Pt(12)
        
        # Use 'content' as a fallback if 'description' or 'summary' is not present
        summary_text = article.get('content', article.get('description', 'No content available.'))
        p = textbox.add_paragraph()
        p.text = summary_text[:200] + "..." if len(summary_text) > 200 else summary_text
        p.font.size = Pt(12)
        textbox.add_paragraph()  # Add spacing

def create_simulation_slide(prs: Presentation, simulation_results: Dict[str, Any]) -> None:
    """Create a slide with simulation results based on simulation_results from SimulationAgent."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Simulation Scenario Analysis"
    
    textbox = slide.shapes.placeholders[1].text_frame
    textbox.clear()
    
    simulations = simulation_results.get('simulation_results', [])
    if not simulations:
        p = textbox.add_paragraph()
        p.text = "No simulation results available."
        p.font.size = Pt(14)
        return

    for sim in simulations:
        p = textbox.add_paragraph()
        p.text = f"📊 Scenario: {sim.get('scenario', 'N/A')}"
        p.font.bold = True
        p.font.size = Pt(14)
        
        p = textbox.add_paragraph()
        p.text = f"Simulated Sales: {sim.get('simulated_sales', 'N/A'):,.2f}"
        p.font.size = Pt(12)
        
        p = textbox.add_paragraph()
        p.text = f"Impact Assessment: {sim.get('impact_assessment', 'N/A')}"
        p.font.size = Pt(12)
        
        p = textbox.add_paragraph()
        p.text = f"Price Elasticity Used: {sim.get('price_elasticity_used', 'N/A'):.2f}"
        p.font.size = Pt(12)
        textbox.add_paragraph()  # Add spacing

class SummaryAgent:
    """Agent to summarize findings and generate reports."""

    def __init__(self, llm_config: Optional[Dict[str, Any]] = None):
        self.llm_config = llm_config

    def llm_generate_summary(self, forecast_result: Dict[str, Any], news_summary: str, simulation_summary: Dict[str, Any], 
                             business_rules: str = None, forecast_trend_analysis: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        prompt = f"""
You are a highly skilled business consultant. Combine the following information into a concise, actionable executive summary and key headlines for a business audience:

Forecast Results:
{json.dumps(forecast_result, indent=2)}

Forecast Trend Analysis:
{json.dumps(forecast_trend_analysis, indent=2) if forecast_trend_analysis else "No significant trends detected."}

News Summary:
{news_summary}

Simulation Summary:
{json.dumps(simulation_summary, indent=2)}

Business Rules/Key Considerations:
{business_rules if business_rules else "None provided."}

Provide:
1. An Executive Summary (2-3 paragraphs)
2. Key Headlines (3-5 bullet points)

Format your response as a JSON object with keys 'executive_summary' and 'headlines'.
"""
        response = call_llm(prompt, max_tokens=1024, llm_config=self.llm_config)
        try:
            return json.loads(response)
        except json.JSONDecodeError:
            logger.error(f"LLM response not valid JSON: {response}")
            return {'executive_summary': "Could not generate summary due to LLM response format.", 'headlines': []}

    def generate_report(self, forecast_result: Dict[str, Any], news_results: Dict[str, Any], simulation_result: Dict[str, Any], user_query: str = "", business_rules: str = ""):
        """Generates a comprehensive report including an LLM summary and a PowerPoint presentation."""
        
        # Generate LLM summary
        summary_data = self.llm_generate_summary(
            forecast_result,
            news_results.get('llm_summary', ''), 
            simulation_result, 
            business_rules,
            forecast_trend_analysis=forecast_result.get('trend_analysis') # Pass trend analysis
        )

        # Generate PowerPoint report
        output_dir = Path('output')
        output_dir.mkdir(exist_ok=True)
        ppt_output_path = output_dir / f'AgenticAI_Report_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'

        # Ensure data formats for PPT are correct
        # forecast_out expects Dict with 'forecast_result' key containing list of dicts
        # news_out expects Dict with 'llm_summary' and 'articles' keys
        # simulation_out expects Dict with 'llm_summary' and 'simulation_results' keys
        # summary_out expects Dict with 'executive_summary' and 'headlines' keys

        generated_ppt_path = create_ppt_report(
            forecast_out=forecast_result,
            news_out=news_results,
            simulation_out=simulation_result,
            summary_out=summary_data,
            output_path=str(ppt_output_path)
        )

        return {
            'summary_data': summary_data,
            'generated_reports': {
                'powerpoint': generated_ppt_path
            }
        }

async def summary_agent_node(state: Dict[str, Any]) -> Dict[str, Any]:
    """Summary agent node for LangGraph workflow, also generates PPT report."""
    logger.info("Executing Summary Agent Node")
    try:
        forecast_result = state.get('forecast_result', {})
        news_results = {
            'articles': state.get('news_results', []),
            'llm_summary': state.get('news_summary', '')
        } # Reconstruct news_results to match expected format
        simulation_result = state.get('simulation_result', {})
        user_query = state.get('query', '')
        business_rules = state.get('business_rules', '') 
        llm_config = state.get('config', {}).get('llm', {})

        agent = SummaryAgent(llm_config=llm_config)
        report_output = agent.generate_report(
            forecast_result, 
            news_results, 
            simulation_result, 
            user_query,
            business_rules
        )

        state['report_output'] = report_output
        return state
        
    except Exception as e:
        logger.error(f"Error in summary agent node: {e}")
        state['error'] = str(e)
        return state