from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import os
from datetime import datetime
import matplotlib.pyplot as plt

def create_ppt_report(forecast_out, news_out, simulation_out, summary_out, output_path='output/AgenticAI_Report.pptx'):
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Agentic AI Business Insights Report"
    slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = summary_out.get('executive_summary', '')

    # Headlines
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Headlines"
    tf = slide.placeholders[1].text_frame
    headlines = summary_out.get('headlines', [])
    for line in headlines:
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(16)

    # Forecast Chart (if available)
    if forecast_out and 'forecast_result' in forecast_out:
        df = pd.DataFrame(forecast_out['forecast_result'])
        if not df.empty and ('yhat' in df or 'yhat_mean' in df):
            plt.figure(figsize=(8,4))
            # Check for both 'date' and 'ds' columns as they might be used by different models
            date_col = 'date' if 'date' in df.columns else 'ds'
            yhat_col = 'yhat' if 'yhat' in df.columns else 'yhat_mean'

            plt.plot(df[date_col], df[yhat_col], label='Forecast')
            
            if 'yhat_lower' in df and 'yhat_upper' in df:
                plt.fill_between(df[date_col], df['yhat_lower'], df['yhat_upper'], alpha=0.2)
            plt.title('Sales Forecast')
            plt.xlabel('Date')
            plt.ylabel('Sales Value')
            plt.legend()
            plt.tight_layout()
            chart_path = 'output/forecast_chart.png'
            plt.savefig(chart_path)
            plt.close()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Forecast Chart"
            left = Inches(1)
            top = Inches(1.5)
            slide.shapes.add_picture(chart_path, left, top, width=Inches(7))
            os.remove(chart_path)

    # Forecast Trend Analysis (if available)
    if forecast_out and 'trend_analysis' in forecast_out:
        trend_analysis = forecast_out['trend_analysis']
        if trend_analysis.get('unusual_trend_detected'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Unusual Forecast Trend Detected"
            tf = slide.placeholders[1].text_frame
            tf.text = trend_analysis.get('trend_description', "")
            
            p = tf.add_paragraph()
            p.text = f"Last 3 forecast values: {trend_analysis.get('last_forecast_values', [])}"
            p.font.size = Pt(14)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Forecast Trend Analysis"
            tf = slide.placeholders[1].text_frame
            tf.text = trend_analysis.get('trend_description', "No significant trends detected.")

    # News Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "News Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = news_out.get('llm_summary', '')
    if news_out.get('articles'):
        for art in news_out['articles'][:5]:
            p = tf.add_paragraph()
            p.text = f"- {art.get('title', 'No Title')}"
            p.font.size = Pt(12)

    # Simulation Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Simulation Insights"
    tf = slide.placeholders[1].text_frame
    tf.text = simulation_out.get('llm_summary', '')
    if simulation_out.get('simulation_results'):
        for sim_res in simulation_out['simulation_results'][:3]:
            p = tf.add_paragraph()
            p.text = f"📊 Scenario: {sim_res.get('scenario', 'N/A')}"
            p.font.bold = True
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = f"Simulated Sales: {sim_res.get('simulated_sales', 'N/A'):.2f}"
            p.font.size = Pt(12)
            
            p = tf.add_paragraph()
            p.text = f"Impact Assessment: {sim_res.get('impact_assessment', 'N/A')}"
            p.font.size = Pt(12)

            if 'price_elasticity_used' in sim_res:
                p = tf.add_paragraph()
                p.text = f"Price Elasticity Used: {sim_res['price_elasticity_used']:.2f}"
                p.font.size = Pt(10)
            textbox.add_paragraph()

    # Save PPT
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path 